"""13-slide professional deck: LLM > Agent (deep) > MCP (deep).
   Expanded agent section + MCP deep dive."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# -- Palette (Dark Theme) --
BG       = RGBColor(0x0B, 0x11, 0x20)
CARD     = RGBColor(0x13, 0x1C, 0x33)
CARD2    = RGBColor(0x1A, 0x25, 0x42)
CYAN     = RGBColor(0x38, 0xBD, 0xF8)
PURPLE   = RGBColor(0xA7, 0x8B, 0xFA)
GREEN    = RGBColor(0x4A, 0xDE, 0x80)
ORANGE   = RGBColor(0xFB, 0x92, 0x3C)
RED      = RGBColor(0xF8, 0x71, 0x71)
YELLOW   = RGBColor(0xFB, 0xD3, 0x4D)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LG       = RGBColor(0xCB, 0xD5, 0xE1)
MG       = RGBColor(0x94, 0xA3, 0xB8)
DG       = RGBColor(0x47, 0x55, 0x69)
TEAL     = RGBColor(0x2D, 0xD4, 0xBF)
PINK     = RGBColor(0xF4, 0x72, 0xB6)
TEXT     = WHITE
CARD_BORDER = RGBColor(0x1A, 0x25, 0x42)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helpers ──

def _bg(s):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG

def _t(s, l, t, w, h, text, sz=18, c=TEXT, b=False, a=PP_ALIGN.LEFT, fn="Segoe UI"):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = b; p.font.name = fn; p.alignment = a
    return tf

def _rect(s, l, t, w, h, f=CARD):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = f; sh.line.fill.background(); sh.shadow.inherit = False
    return sh

def _box(s, l, t, w, h, f=CARD):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = f; sh.line.fill.background(); sh.shadow.inherit = False
    return sh

def _bar(s, l, t, w, c=CYAN):
    return _box(s, l, t, w, Pt(3), c)

def _vbar(s, l, t, h, c=CYAN):
    return _box(s, l, t, Pt(3), h, c)

def _pill(s, l, t, w, h, text, fc=CYAN, tc=WHITE, sz=14):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fc; sh.line.fill.background(); sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.font.color.rgb = tc
    p.font.bold = True; p.font.name = "Segoe UI"; p.alignment = PP_ALIGN.CENTER
    return sh

def _circ(s, l, t, d, fc=CYAN, text="", tc=WHITE, sz=16):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, d, d)
    sh.fill.solid(); sh.fill.fore_color.rgb = fc; sh.line.fill.background()
    if text:
        tf = sh.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.font.color.rgb = tc
        p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    return sh

def _arrow_r(s, l, t, w, c=CYAN):
    sh = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, Inches(0.35))
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background(); sh.shadow.inherit = False

def _arrow_d(s, l, t, h, c=CYAN):
    sh = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, l, t, Inches(0.35), h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background(); sh.shadow.inherit = False

def _code(s, l, t, w, h, text, sz=11, bar_c=CYAN):
    _rect(s, l, t, w, h, RGBColor(0x08, 0x0C, 0x18))
    _box(s, l, t, Pt(3), h, bar_c)
    tb = s.shapes.add_textbox(l + Inches(0.2), t + Inches(0.12), w - Inches(0.35), h - Inches(0.2))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(sz); p.font.name = "Cascadia Code"; p.font.color.rgb = LG; p.space_after = Pt(1)

def _title(s, text, sub="", ac=CYAN):
    _bg(s)
    _box(s, Inches(0), Inches(0), prs.slide_width, Pt(3), ac)
    _t(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), text, sz=32, c=TEXT, b=True)
    _bar(s, Inches(0.8), Inches(1.0), Inches(1.8), ac)
    if sub:
        _t(s, Inches(0.8), Inches(1.15), Inches(11), Inches(0.4), sub, sz=15, c=MG)


# ================================================================
# SLIDE 1 -- Title
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
_circ(s, Inches(-1), Inches(-1), Inches(4), RGBColor(0x10, 0x18, 0x2E))
_circ(s, Inches(10.5), Inches(5), Inches(4), RGBColor(0x10, 0x18, 0x2E))
_bar(s, Inches(0), Inches(0), prs.slide_width, CYAN)

_t(s, Inches(1.5), Inches(1.8), Inches(10.3), Inches(0.5),
   "UNDERSTANDING", sz=18, c=CYAN, b=True, a=PP_ALIGN.CENTER)
_t(s, Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.2),
   "LLMs, Agents & MCP", sz=52, c=TEXT, b=True, a=PP_ALIGN.CENTER)
_bar(s, Inches(5.5), Inches(3.5), Inches(2.3), CYAN)
_t(s, Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.6),
   "A deep dive into AI Agents and the Model Context Protocol", sz=20, c=MG, a=PP_ALIGN.CENTER)

stages = [("LLM", CYAN), ("Agent", PURPLE), ("MCP", GREEN)]
for i, (label, clr) in enumerate(stages):
    x = Inches(3.5) + Inches(i * 2.3)
    _pill(s, x, Inches(5.0), Inches(1.8), Inches(0.6), label, clr, BG, 18)
    if i < 2:
        _arrow_r(s, x + Inches(1.85), Inches(5.1), Inches(0.4), DG)

_t(s, Inches(1.5), Inches(6.4), Inches(10.3), Inches(0.4),
   "Team Knowledge Session  |  2026", sz=13, c=DG, a=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 2 -- What is an LLM
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_title(s, "What is a Large Language Model?", "The AI brain -- powerful but limited", CYAN)

# Input column
_rect(s, Inches(0.5), Inches(1.8), Inches(3.2), Inches(4.8), CARD)
_bar(s, Inches(0.5), Inches(1.8), Inches(3.2), CYAN)
_t(s, Inches(0.7), Inches(1.95), Inches(2.8), Inches(0.35), "INPUT", sz=16, c=CYAN, b=True)
_t(s, Inches(0.7), Inches(2.3), Inches(2.8), Inches(0.3), "Natural language prompt", sz=11, c=MG)
for i, p in enumerate([
    "\"Summarize this document\"", "\"Write Python to sort a list\"",
    "\"What causes rain?\"", "\"Draft a marketing email\"", "\"Review this code for bugs\"",
]):
    _rect(s, Inches(0.7), Inches(2.85) + Inches(i * 0.6), Inches(2.8), Inches(0.45), CARD2)
    _t(s, Inches(0.85), Inches(2.88) + Inches(i * 0.6), Inches(2.5), Inches(0.4), p, sz=11, c=LG)

_arrow_r(s, Inches(3.8), Inches(4.0), Inches(0.6), CYAN)

# LLM Box
_rect(s, Inches(4.5), Inches(1.8), Inches(4.3), Inches(4.8), CARD2)
_bar(s, Inches(4.5), Inches(1.8), Inches(4.3), PURPLE)
_t(s, Inches(4.5), Inches(1.95), Inches(4.3), Inches(0.4), "LARGE LANGUAGE MODEL", sz=16, c=PURPLE, b=True, a=PP_ALIGN.CENTER)
_t(s, Inches(4.5), Inches(2.35), Inches(4.3), Inches(0.3), "GPT-4o  |  Claude  |  Gemini  |  Llama", sz=11, c=MG, a=PP_ALIGN.CENTER)

for i, (lbl, desc) in enumerate([
    ("Trained on", "trillions of words from\nbooks, web, code"),
    ("Contains", "billions of tunable\nparameters (weights)"),
    ("Core task", "predict the most likely\nnext token/word"),
]):
    y = Inches(2.9) + Inches(i * 1.1)
    _t(s, Inches(4.7), y, Inches(1.8), Inches(0.25), lbl, sz=11, c=PURPLE, b=True)
    _t(s, Inches(4.7), y + Inches(0.25), Inches(1.8), Inches(0.7), desc, sz=10, c=LG)

for i, (lbl, desc) in enumerate([
    ("Strength", "understands context,\nnuance, and intent"),
    ("Output", "generates coherent\nhuman-like text"),
    ("Limit", "ONLY produces text.\nCannot act on the world."),
]):
    y = Inches(2.9) + Inches(i * 1.1)
    c2 = RED if i == 2 else PURPLE
    _t(s, Inches(6.7), y, Inches(1.8), Inches(0.25), lbl, sz=11, c=c2, b=True)
    _t(s, Inches(6.7), y + Inches(0.25), Inches(1.8), Inches(0.7), desc, sz=10, c=RED if i == 2 else LG)

_arrow_r(s, Inches(8.9), Inches(4.0), Inches(0.6), PURPLE)

# Output column
_rect(s, Inches(9.6), Inches(1.8), Inches(3.2), Inches(4.8), CARD)
_bar(s, Inches(9.6), Inches(1.8), Inches(3.2), GREEN)
_t(s, Inches(9.8), Inches(1.95), Inches(2.8), Inches(0.35), "OUTPUT", sz=16, c=GREEN, b=True)
_t(s, Inches(9.8), Inches(2.3), Inches(2.8), Inches(0.3), "Generated text only", sz=11, c=MG)
for i, o in enumerate([
    "A 3-paragraph summary", "Working sort function",
    "Scientific explanation", "Polished email draft", "Bug analysis report",
]):
    _rect(s, Inches(9.8), Inches(2.85) + Inches(i * 0.6), Inches(2.8), Inches(0.45), CARD2)
    _t(s, Inches(9.95), Inches(2.88) + Inches(i * 0.6), Inches(2.5), Inches(0.4), o, sz=11, c=GREEN)

_rect(s, Inches(2.0), Inches(6.8), Inches(9.3), Inches(0.5), CARD2)
_t(s, Inches(2.0), Inches(6.82), Inches(9.3), Inches(0.45),
   "KEY LIMITATION:  LLMs generate text but cannot browse the web, read files, call APIs, or take actions.",
   sz=13, c=YELLOW, a=PP_ALIGN.CENTER, b=True)


# ================================================================
# SLIDE 3 -- The Gap: LLMs Can't Act
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_title(s, "The Gap: LLMs Can Think, But Can't Act", "What if we gave them hands?", RED)

# Left: CAN
_rect(s, Inches(0.4), Inches(1.7), Inches(5.8), Inches(4.0), CARD)
_bar(s, Inches(0.4), Inches(1.7), Inches(5.8), GREEN)
_t(s, Inches(0.6), Inches(1.85), Inches(5.0), Inches(0.35), "WHAT LLMs CAN DO", sz=18, c=GREEN, b=True)

for i, (item, detail) in enumerate([
    ("Generate text, code, summaries", "Great at writing, translating, explaining"),
    ("Reason about complex problems", "Logic, math, multi-step analysis"),
    ("Understand context & intent", "Reads between the lines, catches nuance"),
    ("Answer from training knowledge", "Trained on massive datasets up to cutoff"),
]):
    y = Inches(2.4) + Inches(i * 0.7)
    _pill(s, Inches(0.6), y, Inches(0.35), Inches(0.35), "+", GREEN, BG, 14)
    _t(s, Inches(1.1), y, Inches(4.8), Inches(0.3), item, sz=14, c=GREEN, b=True)
    _t(s, Inches(1.1), y + Inches(0.3), Inches(4.8), Inches(0.25), detail, sz=11, c=MG)

# Right: CANNOT
_rect(s, Inches(6.5), Inches(1.7), Inches(6.4), Inches(4.0), CARD)
_bar(s, Inches(6.5), Inches(1.7), Inches(6.4), RED)
_t(s, Inches(6.7), Inches(1.85), Inches(5.0), Inches(0.35), "WHAT LLMs CANNOT DO", sz=18, c=RED, b=True)

for i, (item, detail) in enumerate([
    ("Access live data or the internet", "Training data has a cutoff -- no real-time info"),
    ("Read your files, DBs, or systems", "Cannot see anything outside the prompt"),
    ("Take actions in the real world", "Can't send emails, create tickets, deploy code"),
    ("Execute code or call APIs", "Can write code, but cannot run it"),
]):
    y = Inches(2.4) + Inches(i * 0.7)
    _pill(s, Inches(6.7), y, Inches(0.35), Inches(0.35), "x", RED, WHITE, 14)
    _t(s, Inches(7.2), y, Inches(5.4), Inches(0.3), item, sz=14, c=RED, b=True)
    _t(s, Inches(7.2), y + Inches(0.3), Inches(5.4), Inches(0.25), detail, sz=11, c=MG)

# Bottom: The bridge
_arrow_d(s, Inches(6.3), Inches(4.8), Inches(0.5), ORANGE)
_rect(s, Inches(1.5), Inches(5.5), Inches(10.3), Inches(1.8), CARD)
_bar(s, Inches(1.5), Inches(5.5), Inches(10.3), ORANGE)
_t(s, Inches(1.7), Inches(5.65), Inches(9.8), Inches(0.4),
   "THE SOLUTION: Give the LLM superpowers", sz=20, c=ORANGE, b=True)

powers = [
    ("+ Tools", "Let it call APIs, run code, access databases", CYAN),
    ("+ Memory", "Let it remember context across interactions", GREEN),
    ("+ Autonomy", "Let it plan multi-step tasks and self-correct", PURPLE),
]
for i, (what, desc, clr) in enumerate(powers):
    x = Inches(1.7) + Inches(i * 3.4)
    _pill(s, x, Inches(6.2), Inches(1.2), Inches(0.35), what, clr, BG, 12)
    _t(s, x + Inches(1.3), Inches(6.2), Inches(2.0), Inches(0.35), desc, sz=12, c=LG)

_t(s, Inches(1.5), Inches(6.85), Inches(10.3), Inches(0.35),
   "LLM + Tools + Memory + Autonomy  =  AGENT", sz=18, c=YELLOW, b=True, a=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 4 -- What is an Agent (Architecture)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_title(s, "What is an AI Agent?", "An autonomous system that perceives, reasons, plans, and acts", PURPLE)

# Center: LLM brain
_rect(s, Inches(4.2), Inches(2.2), Inches(4.8), Inches(2.2), CARD2)
_bar(s, Inches(4.2), Inches(2.2), Inches(4.8), PURPLE)
_t(s, Inches(4.2), Inches(2.3), Inches(4.8), Inches(0.4), "LLM  (The Reasoning Engine)", sz=18, c=PURPLE, b=True, a=PP_ALIGN.CENTER)
_t(s, Inches(4.2), Inches(2.8), Inches(4.8), Inches(0.3),
   "Understands intent  |  Makes decisions  |  Self-corrects", sz=12, c=LG, a=PP_ALIGN.CENTER)
# Core capabilities inside
for i, (lbl, clr) in enumerate([("Perceive", CYAN), ("Reason", PURPLE), ("Plan", GREEN), ("Act", ORANGE)]):
    x = Inches(4.5) + Inches(i * 1.1)
    _pill(s, x, Inches(3.3), Inches(0.95), Inches(0.35), lbl, clr, BG, 11)

# Left: TOOLS
_rect(s, Inches(0.3), Inches(1.7), Inches(3.3), Inches(4.2), CARD)
_bar(s, Inches(0.3), Inches(1.7), Inches(3.3), CYAN)
_t(s, Inches(0.5), Inches(1.85), Inches(2.8), Inches(0.35), "TOOLS", sz=18, c=CYAN, b=True)
_t(s, Inches(0.5), Inches(2.2), Inches(2.8), Inches(0.25), "What the agent can DO", sz=11, c=MG)
for i, t in enumerate(["Call APIs & web services", "Execute code & commands", "Read & write files", "Query databases", "Search the internet", "Send messages & emails"]):
    _pill(s, Inches(0.5), Inches(2.6) + Inches(i * 0.42), Inches(2.8), Inches(0.32), t, CARD2, LG, 10)
_arrow_r(s, Inches(3.65), Inches(3.1), Inches(0.5), CYAN)

# Right: MEMORY
_rect(s, Inches(9.7), Inches(1.7), Inches(3.3), Inches(4.2), CARD)
_bar(s, Inches(9.7), Inches(1.7), Inches(3.3), GREEN)
_t(s, Inches(9.9), Inches(1.85), Inches(2.8), Inches(0.35), "MEMORY", sz=18, c=GREEN, b=True)
_t(s, Inches(9.9), Inches(2.2), Inches(2.8), Inches(0.25), "What the agent KNOWS", sz=11, c=MG)
for i, t in enumerate(["Conversation history", "Prior decisions & results", "Retrieved knowledge (RAG)", "User preferences", "Task progress & state", "Long-term learned facts"]):
    _pill(s, Inches(9.9), Inches(2.6) + Inches(i * 0.42), Inches(2.8), Inches(0.32), t, CARD2, LG, 10)
_arrow_r(s, Inches(9.1), Inches(3.1), Inches(0.5), GREEN)

# Bottom: ENVIRONMENT
_rect(s, Inches(0.3), Inches(6.1), Inches(12.7), Inches(1.2), CARD)
_bar(s, Inches(0.3), Inches(6.1), Inches(12.7), ORANGE)
_t(s, Inches(0.5), Inches(6.15), Inches(2.0), Inches(0.3), "ENVIRONMENT", sz=14, c=ORANGE, b=True)
env_items = ["File system", "Databases", "APIs", "Web browsers", "Terminals", "Cloud services", "Users"]
for i, item in enumerate(env_items):
    x = Inches(2.5) + Inches(i * 1.5)
    _pill(s, x, Inches(6.5), Inches(1.3), Inches(0.4), item, CARD2, LG, 10)

# Equation
_rect(s, Inches(3.5), Inches(4.7), Inches(6.3), Inches(0.5), CARD2)
_t(s, Inches(3.5), Inches(4.72), Inches(6.3), Inches(0.45),
   "AGENT  =  LLM  +  Tools  +  Memory  +  Environment", sz=16, c=YELLOW, b=True, a=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 5 -- How Agents Work: The Loop (detailed)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_title(s, "How Agents Work: The Agent Loop", "The core pattern behind every AI agent", PURPLE)

# Big circular loop diagram
# We'll show it as a detailed horizontal flow with a loop-back

# Main flow cards
loop_steps = [
    ("1", "PERCEIVE", "Agent receives input:\nuser request, tool results,\nor environment changes",
     "What happened?", CYAN),
    ("2", "REASON", "LLM analyzes the situation:\nwhat's the goal, what do\nI know, what's missing?",
     "What do I know?", PURPLE),
    ("3", "PLAN", "Break the task into steps:\nprioritize, identify tools\nneeded, estimate approach",
     "What should I do?", GREEN),
    ("4", "ACT", "Execute the next step:\ncall a tool, write code,\nquery a database, etc.",
     "Do the thing!", ORANGE),
    ("5", "OBSERVE", "Read the result of the\naction: did it work?\nany errors? new data?",
     "What happened?", TEAL),
    ("6", "REFLECT", "Evaluate progress: am I\ndone? do I need to retry?\nshould I change my plan?",
     "How did it go?", PINK),
]

for i, (num, title, desc, question, clr) in enumerate(loop_steps):
    x = Inches(0.3) + Inches(i * 2.15)
    _rect(s, x, Inches(1.7), Inches(2.0), Inches(3.5), CARD)
    _bar(s, x, Inches(1.7), Inches(2.0), clr)
    _circ(s, x + Inches(0.1), Inches(1.8), Inches(0.4), clr, num, BG, 14)
    _t(s, x + Inches(0.6), Inches(1.82), Inches(1.3), Inches(0.3), title, sz=14, c=clr, b=True)
    _t(s, x + Inches(0.1), Inches(2.3), Inches(1.8), Inches(1.5), desc, sz=10, c=LG)
    # Question badge
    _pill(s, x + Inches(0.1), Inches(4.0), Inches(1.8), Inches(0.3), question, CARD2, clr, 9)
    # Arrow to next
    if i < 5:
        _arrow_r(s, x + Inches(2.0), Inches(2.9), Inches(0.15), DG)

# Loop-back arrow (visual)
_rect(s, Inches(0.3), Inches(5.4), Inches(12.7), Inches(0.6), CARD2)
_t(s, Inches(0.3), Inches(5.42), Inches(12.7), Inches(0.55),
   "LOOP BACK to step 1 if not done  ---  Each iteration gets closer to the goal  ---  Agent decides when to stop",
   sz=13, c=YELLOW, a=PP_ALIGN.CENTER, b=True)

# Bottom: Key insight
_rect(s, Inches(0.3), Inches(6.2), Inches(12.7), Inches(1.0), CARD)
_t(s, Inches(0.5), Inches(6.25), Inches(4.0), Inches(0.3), "WHY THIS MATTERS", sz=14, c=TEXT, b=True)

insights = [
    ("Self-correcting", "If a tool call fails, agent retries with a different approach", CYAN),
    ("Multi-step", "Complex tasks are broken down and executed incrementally", GREEN),
    ("Autonomous", "No human needed between steps -- agent drives itself", ORANGE),
]
for i, (lbl, desc, clr) in enumerate(insights):
    x = Inches(0.5) + Inches(i * 4.2)
    _pill(s, x, Inches(6.65), Inches(1.5), Inches(0.3), lbl, clr, BG, 10)
    _t(s, x + Inches(1.6), Inches(6.65), Inches(2.5), Inches(0.3), desc, sz=10, c=LG)


# ================================================================
# SLIDE 6 -- 4 Agentic Design Patterns
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_title(s, "4 Agentic Design Patterns", "Building blocks you can combine for any AI workflow", ORANGE)

patterns = [
    ("REFLECTION", CYAN,
     "Agent reviews its own output\nbefore delivering.",
     ["Generate initial output",
      "Critique: \"Is this correct?\"",
      "Find gaps or errors",
      "Revise and improve",
      "Deliver polished result"],
     "Code review, writing, QA",
     "\"Write tests, then review\nthem for edge cases\""),

    ("TOOL USE", PURPLE,
     "Agent decides which tools\nto call and when.",
     ["Analyze what's needed",
      "Select the right tool",
      "Call tool with parameters",
      "Process the result",
      "Decide: done or next tool?"],
     "Data analysis, automation",
     "\"Look up the customer,\nthen check their orders\""),

    ("PLANNING", GREEN,
     "Agent breaks complex tasks\ninto ordered sub-tasks.",
     ["Decompose into sub-tasks",
      "Order by dependencies",
      "Execute step by step",
      "Re-plan if something fails",
      "Verify all steps complete"],
     "Research, project scoping",
     "\"Migrate this API:\nplan steps first\""),

    ("MULTI-AGENT", ORANGE,
     "Specialized agents collaborate\non different parts.",
     ["Planner agent designs",
      "Coder agent implements",
      "Reviewer agent checks",
      "Agents hand off work",
      "Orchestrator manages all"],
     "Software dev, content",
     "\"Architect + Developer +\nQA each do their part\""),
]

for i, (title, clr, desc, steps, used_in, example) in enumerate(patterns):
    x = Inches(0.2) + Inches(i * 3.3)
    _rect(s, x, Inches(1.6), Inches(3.1), Inches(5.7), CARD)
    _bar(s, x, Inches(1.6), Inches(3.1), clr)

    # Title
    _t(s, x + Inches(0.15), Inches(1.7), Inches(2.8), Inches(0.3), title, sz=16, c=clr, b=True)
    # Description
    _t(s, x + Inches(0.15), Inches(2.05), Inches(2.8), Inches(0.7), desc, sz=10, c=MG)

    # Numbered steps with connectors
    for j, step in enumerate(steps):
        sy = Inches(2.85) + Inches(j * 0.5)
        _circ(s, x + Inches(0.15), sy + Inches(0.02), Inches(0.25), clr, str(j+1), BG, 9)
        _t(s, x + Inches(0.5), sy, Inches(2.4), Inches(0.25), step, sz=10, c=LG)
        if j < 4:
            _box(s, x + Inches(0.26), sy + Inches(0.27), Pt(2), Inches(0.25), DG)

    # Example box
    _rect(s, x + Inches(0.1), Inches(5.55), Inches(2.85), Inches(0.95), CARD2)
    _t(s, x + Inches(0.2), Inches(5.58), Inches(2.6), Inches(0.2), "Example:", sz=9, c=clr, b=True)
    _t(s, x + Inches(0.2), Inches(5.8), Inches(2.6), Inches(0.6), example, sz=9, c=LG)

    # Used in
    _t(s, x + Inches(0.15), Inches(6.55), Inches(2.8), Inches(0.2), used_in, sz=9, c=MG)


# ================================================================
# SLIDE 7 -- Agent in Action (before MCP, show raw tool use)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_title(s, "Agent in Action: A Real Example",
       "Watch an agent handle: \"Find the bug in prod, fix it, and notify the team\"", PURPLE)

# Left: step-by-step narration
steps_data = [
    ("1", "PERCEIVE", "User asks: \"There's a 500 error in prod, fix it\"", CYAN),
    ("2", "PLAN", "Agent decides: check logs > find error > read code > fix > test > notify", PURPLE),
    ("3", "ACT: Search logs", "Calls log search tool: search_logs(status=500, last='1h')", GREEN),
    ("4", "OBSERVE", "Gets result: NullPointerError in /api/users line 42", TEAL),
    ("5", "ACT: Read code", "Calls file read tool: read_file('api/users.py', line=42)", GREEN),
    ("6", "REASON", "Identifies: user.email accessed before null check", PURPLE),
    ("7", "ACT: Fix code", "Calls edit tool: edit_file('api/users.py', ...) -- adds null check", ORANGE),
    ("8", "ACT: Run tests", "Calls test tool: run_tests('test_users.py') -- all pass!", GREEN),
    ("9", "ACT: Notify", "Calls Slack tool: send_message('#dev', 'Fixed NullPtr in /api/users')", PINK),
    ("10", "DONE", "Reports to user: \"Fixed! Added null check, tests pass, team notified.\"", YELLOW),
]

for i, (num, phase, desc, clr) in enumerate(steps_data):
    y = Inches(1.5) + Inches(i * 0.58)
    _circ(s, Inches(0.4), y + Inches(0.04), Inches(0.38), clr, num, BG, 11)
    if i < 9:
        _box(s, Inches(0.57), y + Inches(0.42), Pt(2), Inches(0.2), DG)
    _pill(s, Inches(0.9), y, Inches(1.3), Inches(0.32), phase, CARD2, clr, 9)
    _t(s, Inches(2.3), y + Inches(0.02), Inches(5.0), Inches(0.3), desc, sz=11, c=LG)

# Right: Visual showing tools used
_rect(s, Inches(7.5), Inches(1.5), Inches(5.5), Inches(5.7), CARD)
_t(s, Inches(7.7), Inches(1.6), Inches(5.0), Inches(0.3), "TOOLS THE AGENT USED", sz=14, c=TEXT, b=True)

tools_used = [
    ("search_logs()", "Find errors in production logs", CYAN, "Log system"),
    ("read_file()", "Read source code files", GREEN, "File system"),
    ("edit_file()", "Modify code to fix the bug", ORANGE, "File system"),
    ("run_tests()", "Execute test suite", PURPLE, "Terminal"),
    ("send_message()", "Notify team on Slack", PINK, "Slack API"),
]
for i, (func, desc, clr, source) in enumerate(tools_used):
    y = Inches(2.1) + Inches(i * 0.75)
    _rect(s, Inches(7.7), y, Inches(5.0), Inches(0.6), CARD2)
    _vbar(s, Inches(7.7), y, Inches(0.6), clr)
    _t(s, Inches(7.95), y + Inches(0.02), Inches(2.2), Inches(0.25), func, sz=12, c=clr, b=True, fn="Cascadia Code")
    _t(s, Inches(7.95), y + Inches(0.27), Inches(2.5), Inches(0.25), desc, sz=10, c=LG)
    _pill(s, Inches(10.8), y + Inches(0.1), Inches(1.7), Inches(0.3), source, CARD, MG, 9)

# Key insight
_rect(s, Inches(7.7), Inches(5.9), Inches(5.0), Inches(1.1), CARD2)
_t(s, Inches(7.9), Inches(5.95), Inches(4.6), Inches(0.3),
   "THE PROBLEM", sz=13, c=RED, b=True)
_t(s, Inches(7.9), Inches(6.3), Inches(4.6), Inches(0.6),
   "Each tool is a custom integration.\nWhat if there was a standard way to connect\nALL these tools?  Enter MCP...", sz=12, c=LG)


# ================================================================
# SLIDE 8 -- What is MCP (deep)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_title(s, "What is MCP?", "Model Context Protocol -- the universal connector for AI", GREEN)

# Analogy
_rect(s, Inches(0.4), Inches(1.7), Inches(5.8), Inches(2.6), CARD)
_bar(s, Inches(0.4), Inches(1.7), Inches(5.8), GREEN)
_t(s, Inches(0.6), Inches(1.85), Inches(5.4), Inches(0.35), "THE ANALOGY", sz=16, c=GREEN, b=True)

_t(s, Inches(0.6), Inches(2.35), Inches(1.4), Inches(0.3), "Before USB-C:", sz=13, c=MG, b=True)
_t(s, Inches(2.0), Inches(2.35), Inches(4.0), Inches(0.3),
   "Every device had its own charger. Micro-USB, Lightning, barrel jacks...", sz=13, c=LG)
_t(s, Inches(0.6), Inches(2.85), Inches(1.4), Inches(0.3), "After USB-C:", sz=13, c=GREEN, b=True)
_t(s, Inches(2.0), Inches(2.85), Inches(4.0), Inches(0.3),
   "One universal port. Any device, any charger. Just works.", sz=13, c=GREEN)
_t(s, Inches(0.6), Inches(3.45), Inches(2.5), Inches(0.3), "MCP is USB-C for AI.", sz=16, c=YELLOW, b=True)
_t(s, Inches(3.1), Inches(3.45), Inches(3.0), Inches(0.3),
   "One protocol. Any AI app. Any tool.", sz=14, c=YELLOW)

# Technical
_rect(s, Inches(6.5), Inches(1.7), Inches(6.4), Inches(2.6), CARD)
_bar(s, Inches(6.5), Inches(1.7), Inches(6.4), CYAN)
_t(s, Inches(6.7), Inches(1.85), Inches(5.0), Inches(0.35), "TECHNICALLY", sz=16, c=CYAN, b=True)

for i, (lbl, desc) in enumerate([
    ("Open standard", "by Anthropic, adopted industry-wide (GitHub, Google, etc.)"),
    ("JSON-RPC 2.0", "lightweight request/response protocol"),
    ("Transport", "stdio (local processes) or HTTP+SSE (remote servers)"),
    ("Discovery", "AI auto-discovers available tools at startup"),
    ("Stateful", "maintains session context between AI and tools"),
]):
    y = Inches(2.35) + Inches(i * 0.35)
    _t(s, Inches(6.9), y, Inches(1.8), Inches(0.3), lbl, sz=11, c=CYAN, b=True)
    _t(s, Inches(8.8), y, Inches(3.8), Inches(0.3), desc, sz=11, c=LG)

# N x M vs N + M
_rect(s, Inches(0.4), Inches(4.6), Inches(5.8), Inches(2.7), CARD)
_t(s, Inches(0.6), Inches(4.7), Inches(3.0), Inches(0.35), "WITHOUT MCP", sz=14, c=RED, b=True)
_t(s, Inches(3.5), Inches(4.75), Inches(2.5), Inches(0.3), "N x M integrations", sz=12, c=RED)

for i, app in enumerate(["Copilot", "Cursor", "Custom"]):
    _pill(s, Inches(0.5), Inches(5.2) + Inches(i * 0.55), Inches(1.3), Inches(0.4), app, CARD2, MG, 10)
for i, tool in enumerate(["GitHub", "Slack", "DB"]):
    _pill(s, Inches(4.2), Inches(5.2) + Inches(i * 0.55), Inches(1.3), Inches(0.4), tool, CARD2, MG, 10)
for ai in range(3):
    y = Inches(5.37) + Inches(ai * 0.55)
    _box(s, Inches(1.85), y, Inches(2.35), Pt(1), RGBColor(0x50, 0x25, 0x25))

_rect(s, Inches(6.5), Inches(4.6), Inches(6.4), Inches(2.7), CARD)
_t(s, Inches(6.7), Inches(4.7), Inches(3.0), Inches(0.35), "WITH MCP", sz=14, c=GREEN, b=True)
_t(s, Inches(9.5), Inches(4.75), Inches(3.0), Inches(0.3), "N + M integrations", sz=12, c=GREEN)

for i, app in enumerate(["Copilot", "Cursor", "Custom"]):
    _pill(s, Inches(6.6), Inches(5.2) + Inches(i * 0.55), Inches(1.3), Inches(0.4), app, CARD2, CYAN, 10)
    _arrow_r(s, Inches(7.95), Inches(5.25) + Inches(i * 0.55), Inches(0.35), CYAN)

_rect(s, Inches(8.5), Inches(5.1), Inches(1.2), Inches(2.0), GREEN)
_t(s, Inches(8.5), Inches(5.85), Inches(1.2), Inches(0.4), "MCP", sz=16, c=WHITE, b=True, a=PP_ALIGN.CENTER)  # MCP hub text

for i, tool in enumerate(["GitHub", "Slack", "DB"]):
    _arrow_r(s, Inches(9.75), Inches(5.25) + Inches(i * 0.55), Inches(0.35), GREEN)
    _pill(s, Inches(10.2), Inches(5.2) + Inches(i * 0.55), Inches(1.3), Inches(0.4), tool, CARD2, GREEN, 10)


# ================================================================
# SLIDE 9 -- MCP Architecture (detailed)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_title(s, "MCP Architecture", "How Host, Client, Server, and Primitives work together", CYAN)

_rect(s, Inches(0.3), Inches(1.6), Inches(12.7), Inches(4.2), RGBColor(0x0E, 0x15, 0x28))
_t(s, Inches(0.5), Inches(1.65), Inches(5.0), Inches(0.3), "MCP HOST", sz=13, c=CYAN, b=True)
_t(s, Inches(2.0), Inches(1.65), Inches(6.0), Inches(0.3), "(VS Code, Copilot, Cursor, or your custom app)", sz=11, c=MG)

# LLM
_rect(s, Inches(0.5), Inches(2.1), Inches(2.8), Inches(3.4), CARD)
_bar(s, Inches(0.5), Inches(2.1), Inches(2.8), PURPLE)
_t(s, Inches(0.5), Inches(2.2), Inches(2.8), Inches(0.35), "LLM", sz=18, c=PURPLE, b=True, a=PP_ALIGN.CENTER)
_t(s, Inches(0.5), Inches(2.6), Inches(2.8), Inches(0.3), "GPT-4o / Claude / Gemini", sz=10, c=MG, a=PP_ALIGN.CENTER)
_t(s, Inches(0.7), Inches(3.1), Inches(2.4), Inches(0.25), "Sees available tools", sz=11, c=LG)
_t(s, Inches(0.7), Inches(3.4), Inches(2.4), Inches(0.25), "Decides which to call", sz=11, c=LG)
_t(s, Inches(0.7), Inches(3.7), Inches(2.4), Inches(0.25), "Processes results", sz=11, c=LG)
_t(s, Inches(0.7), Inches(4.0), Inches(2.4), Inches(0.25), "Generates response", sz=11, c=LG)

_arrow_r(s, Inches(3.4), Inches(3.5), Inches(0.55), PURPLE)

# Clients
clients_y = [Inches(2.1), Inches(3.45), Inches(4.8)]
client_colors = [GREEN, ORANGE, PINK]
for i, (y, clr) in enumerate(zip(clients_y, client_colors)):
    _rect(s, Inches(4.1), y, Inches(1.7), Inches(1.0), CARD2)
    _bar(s, Inches(4.1), y, Inches(1.7), clr)
    _t(s, Inches(4.1), y + Inches(0.08), Inches(1.7), Inches(0.3), f"Client {i+1}", sz=12, c=clr, b=True, a=PP_ALIGN.CENTER)
    _t(s, Inches(4.1), y + Inches(0.4), Inches(1.7), Inches(0.5), "1:1 session\nJSON-RPC 2.0", sz=9, c=MG, a=PP_ALIGN.CENTER)
    _arrow_r(s, Inches(5.9), y + Inches(0.3), Inches(0.5), clr)

# Servers
server_data = [
    ("GitHub MCP Server", ["search_code()", "create_pr()", "list_issues()"], GREEN),
    ("Database Server", ["query(sql)", "get_schema()", "insert_row()"], ORANGE),
    ("Custom Server", ["deploy(env)", "check_status()", "rollback()"], PINK),
]
for i, (name, tools, clr) in enumerate(server_data):
    y = clients_y[i]
    _rect(s, Inches(6.5), y, Inches(2.8), Inches(1.0), CARD)
    _bar(s, Inches(6.5), y, Inches(2.8), clr)
    _t(s, Inches(6.6), y + Inches(0.05), Inches(2.6), Inches(0.25), name, sz=11, c=clr, b=True)
    _t(s, Inches(6.6), y + Inches(0.35), Inches(2.6), Inches(0.55), "   ".join(tools), sz=8, c=LG)
    _arrow_r(s, Inches(9.4), y + Inches(0.3), Inches(0.45), clr)

ext = [("GitHub API", GREEN), ("PostgreSQL", ORANGE), ("AWS / Cloud", PINK)]
for i, (name, clr) in enumerate(ext):
    _pill(s, Inches(10.0), clients_y[i] + Inches(0.1), Inches(2.0), Inches(0.7), name, clr, BG, 11)

_t(s, Inches(5.95), Inches(1.7), Inches(0.5), Inches(0.25), "stdio", sz=9, c=MG, a=PP_ALIGN.CENTER)
_t(s, Inches(5.9), Inches(4.3), Inches(0.6), Inches(0.25), "HTTP", sz=9, c=MG, a=PP_ALIGN.CENTER)

# Primitives bar
_rect(s, Inches(0.3), Inches(6.0), Inches(12.7), Inches(1.3), CARD)
_t(s, Inches(0.5), Inches(6.05), Inches(2.0), Inches(0.3), "MCP PRIMITIVES", sz=13, c=TEXT, b=True)

prims = [
    ("TOOLS", "Functions AI calls to perform actions", "create_issue(), send_msg(), deploy()", CYAN),
    ("RESOURCES", "Data AI reads for context (read-only)", "file contents, DB schemas, configs", GREEN),
    ("PROMPTS", "Reusable instruction templates", "\"summarize-pr\", \"triage-bug\"", PURPLE),
]
for i, (name, desc, ex, clr) in enumerate(prims):
    x = Inches(0.5) + Inches(i * 4.2)
    _pill(s, x, Inches(6.4), Inches(1.2), Inches(0.3), name, clr, BG, 11)
    _t(s, x + Inches(1.3), Inches(6.38), Inches(2.5), Inches(0.25), desc, sz=11, c=TEXT, b=True)
    _t(s, x + Inches(1.3), Inches(6.65), Inches(2.7), Inches(0.25), ex, sz=9, c=MG)


# ================================================================
# SLIDE 10 -- What Can Become an MCP Server
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_title(s, "What Can Become an MCP Server?", "If it has an API, a CLI, or data -- it can be an MCP server", GREEN)

categories = [
    ("Dev Tools", GREEN,
     ["GitHub / GitLab", "Jira / Linear", "CI/CD pipelines", "Docker / K8s", "Terraform"],
     "Any tool your devs use daily"),
    ("Databases", ORANGE,
     ["PostgreSQL / MySQL", "MongoDB / DynamoDB", "Redis / Elasticsearch", "S3 / Blob storage", "Data warehouses"],
     "Any data source you query"),
    ("Communication", PURPLE,
     ["Slack / Teams", "Email (SMTP/Graph)", "Confluence / Notion", "Calendar / Scheduling", "PagerDuty / Opsgenie"],
     "Any channel you communicate on"),
    ("Cloud & Infra", CYAN,
     ["AWS services", "Azure / GCP APIs", "Monitoring (Datadog)", "Log systems (Splunk)", "DNS / CDN / LBs"],
     "Any infra you manage"),
    ("Internal Tools", PINK,
     ["Admin dashboards", "CRM (Salesforce)", "HR systems", "Billing / Payments", "Custom REST APIs"],
     "Any internal system with an API"),
    ("Knowledge", TEAL,
     ["Documentation sites", "Wikis / Runbooks", "Code repositories", "PDF / Office docs", "Vector DBs (RAG)"],
     "Any knowledge your team needs"),
]

for i, (title, clr, items, tagline) in enumerate(categories):
    col = i % 3
    row = i // 3
    x = Inches(0.3) + Inches(col * 4.3)
    y = Inches(1.6) + Inches(row * 2.85)
    _rect(s, x, y, Inches(4.1), Inches(2.65), CARD)
    _bar(s, x, y, Inches(4.1), clr)
    _t(s, x + Inches(0.2), y + Inches(0.1), Inches(3.6), Inches(0.3), title, sz=16, c=clr, b=True)
    _t(s, x + Inches(0.2), y + Inches(0.4), Inches(3.6), Inches(0.25), tagline, sz=10, c=MG)
    for j, item in enumerate(items):
        _pill(s, x + Inches(0.15), y + Inches(0.8) + Inches(j * 0.35), Inches(3.7), Inches(0.28), item, CARD2, LG, 10)


# ================================================================
# SLIDE 11 -- How to Convert + Build MCP Server
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_title(s, "Build an MCP Server: From Idea to Working Code", "5-step pattern + real Jira example your team can use today", ORANGE)

# Left: The 5-step pattern
steps = [
    ("1", "IDENTIFY", "Pick a tool your team uses that has an API or CLI", CYAN),
    ("2", "MAP", "List the top actions (tools) and data reads (resources)", PURPLE),
    ("3", "CODE", "Write a FastMCP server -- thin wrapper around the existing API", GREEN),
    ("4", "CONFIG", "Add to VS Code settings.json -- Copilot auto-discovers it", ORANGE),
    ("5", "USE", "Ask Copilot in natural language -- it calls your tools!", YELLOW),
]
for i, (num, title, desc, clr) in enumerate(steps):
    y = Inches(1.55) + Inches(i * 0.7)
    _circ(s, Inches(0.4), y + Inches(0.05), Inches(0.4), clr, num, BG, 14)
    if i < 4:
        _box(s, Inches(0.58), y + Inches(0.45), Pt(2), Inches(0.3), DG)
    _t(s, Inches(0.95), y + Inches(0.02), Inches(1.0), Inches(0.3), title, sz=13, c=clr, b=True)
    _t(s, Inches(2.0), y + Inches(0.02), Inches(3.8), Inches(0.3), desc, sz=12, c=LG)

# Right: Real code
code_jira = '''# jira_server.py
from mcp.server.fastmcp import FastMCP
from jira import JIRA

mcp = FastMCP("Jira Tools")
jira = JIRA(server="https://team.atlassian.net",
            basic_auth=("email", "token"))

@mcp.tool()
def get_my_tickets(status: str = "In Progress"):
    """Get tickets assigned to me."""
    issues = jira.search_issues(
      f'assignee=currentUser() AND status="{status}"')
    return "\\n".join(
      f"{i.key}: {i.fields.summary}" for i in issues)

@mcp.tool()
def create_ticket(project: str, summary: str):
    """Create a new Jira ticket."""
    issue = jira.create_issue(
      project=project, summary=summary,
      issuetype={"name": "Task"})
    return f"Created {issue.key}"

@mcp.resource("jira://current-sprint")
def current_sprint():
    """Get all tickets in current sprint."""
    issues = jira.search_issues(
      'sprint in openSprints()')
    return "\\n".join(
      f"{i.key} [{i.fields.status}] {i.fields.summary}"
      for i in issues)

if __name__ == "__main__":
    mcp.run(transport="stdio")'''

_code(s, Inches(5.8), Inches(1.5), Inches(5.0), Inches(4.3), code_jira, sz=9, bar_c=GREEN)

# VS Code config
vs_code = '''// .vscode/settings.json
{ "mcp": { "servers": { "jira": {
    "command": "python",
    "args": ["jira_server.py"]
}}}}'''
_code(s, Inches(5.8), Inches(5.95), Inches(5.0), Inches(0.9), vs_code, sz=9, bar_c=CYAN)

# Bottom left: "Now ask Copilot"
_rect(s, Inches(0.3), Inches(5.2), Inches(5.2), Inches(2.1), CARD)
_bar(s, Inches(0.3), Inches(5.2), Inches(5.2), GREEN)
_t(s, Inches(0.5), Inches(5.3), Inches(4.8), Inches(0.3), "NOW ASK COPILOT:", sz=14, c=GREEN, b=True)

for i, p in enumerate([
    "\"Show my in-progress tickets\"",
    "\"Create a bug for the login crash\"",
    "\"What's in our current sprint?\"",
    "\"Create tickets for each TODO in this file\"",
]):
    _rect(s, Inches(0.5), Inches(5.75) + Inches(i * 0.38), Inches(4.8), Inches(0.3), CARD2)
    _t(s, Inches(0.65), Inches(5.77) + Inches(i * 0.38), Inches(4.5), Inches(0.28), p, sz=11, c=GREEN)

_t(s, Inches(10.85), Inches(6.95), Inches(2.0), Inches(0.25),
   "~50 lines of Python!", sz=10, c=YELLOW, b=True)


# ================================================================
# SLIDE 12 -- Agent + MCP End-to-End
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_title(s, "End-to-End: Agent + MCP in Action",
       "Real scenario: \"Triage all new bugs and assign them to the right developer\"", ORANGE)

flow = [
    ("1", "You ask Copilot", "\"Triage new bugs: check severity, assign to right dev, notify on Slack\"", CYAN),
    ("2", "Agent plans", "Sub-tasks: (a) fetch new bugs (b) analyze each (c) assign (d) notify", PURPLE),
    ("3", "Calls Jira MCP", "get_my_tickets(status='New') -- returns 8 unassigned bugs", GREEN),
    ("4", "LLM analyzes", "Reads title + description, infers severity & component owner", PURPLE),
    ("5", "Calls Jira MCP", "transition_ticket('BUG-42', 'In Progress') for each bug", ORANGE),
    ("6", "Calls Slack MCP", "send_message(channel='#dev', text='8 bugs triaged & assigned')", PINK),
    ("7", "Reports back", "\"Done! 8 bugs triaged: 2 critical (Alice), 6 normal. Team notified.\"", YELLOW),
]

for i, (num, label, desc, clr) in enumerate(flow):
    y = Inches(1.55) + Inches(i * 0.8)
    _circ(s, Inches(0.4), y + Inches(0.05), Inches(0.45), clr, num, BG, 15)
    if i < 6:
        _box(s, Inches(0.6), y + Inches(0.5), Pt(2), Inches(0.35), DG)
    _t(s, Inches(1.0), y, Inches(2.0), Inches(0.3), label, sz=14, c=clr, b=True)
    _t(s, Inches(1.0), y + Inches(0.3), Inches(5.5), Inches(0.3), desc, sz=11, c=LG)

# Right: Architecture
_rect(s, Inches(7.2), Inches(1.5), Inches(5.8), Inches(5.8), CARD)
_t(s, Inches(7.4), Inches(1.6), Inches(5.4), Inches(0.3), "UNDER THE HOOD", sz=13, c=TEXT, b=True)

arch_items = [
    ("YOU", "Type in Copilot Chat", CYAN, Inches(2.1)),
    ("COPILOT AGENT", "LLM reasons + plans steps", PURPLE, Inches(3.0)),
    ("MCP CLIENT", "Routes to the right server", CYAN, Inches(3.9)),
]
for label, desc, clr, y in arch_items:
    _pill(s, Inches(7.5), y, Inches(5.2), Inches(0.6), f"{label}:  {desc}", CARD2, clr, 12)

_t(s, Inches(9.8), Inches(4.55), Inches(0.5), Inches(0.3), "v", sz=20, c=CYAN, a=PP_ALIGN.CENTER)

servers = [("Jira\nServer", GREEN), ("Slack\nServer", PINK), ("GitHub\nServer", ORANGE)]
for i, (name, clr) in enumerate(servers):
    x = Inches(7.5) + Inches(i * 1.85)
    _pill(s, x, Inches(4.9), Inches(1.6), Inches(0.7), name, clr, BG, 11)
    _arrow_d(s, x + Inches(0.6), Inches(5.65), Inches(0.35), clr)

for i, (name, clr) in enumerate([("Jira API", GREEN), ("Slack API", PINK), ("GitHub API", ORANGE)]):
    x = Inches(7.5) + Inches(i * 1.85)
    _pill(s, x, Inches(6.1), Inches(1.6), Inches(0.5), name, CARD2, clr, 10)

_rect(s, Inches(7.4), Inches(6.8), Inches(5.2), Inches(0.4), CARD2)
_t(s, Inches(7.4), Inches(6.82), Inches(5.2), Inches(0.35),
   "Each MCP server is independent & reusable across apps", sz=11, c=YELLOW, b=True, a=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 13 -- Next Steps & Resources
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_title(s, "Start Building Today", "Your action plan and resources", YELLOW)

# Left: Action plan
_rect(s, Inches(0.4), Inches(1.6), Inches(6.3), Inches(5.6), CARD)
_bar(s, Inches(0.4), Inches(1.6), Inches(6.3), YELLOW)
_t(s, Inches(0.6), Inches(1.7), Inches(5.8), Inches(0.35), "YOUR ACTION PLAN", sz=18, c=YELLOW, b=True)

plan = [
    ("THIS WEEK", [
        "Install MCP Python SDK:  pip install mcp",
        "Build a hello-world MCP server (15 min)",
        "Connect it to VS Code Copilot",
    ], CYAN),
    ("NEXT WEEK", [
        "Pick ONE internal tool your team uses daily",
        "Map its top 3-5 operations to MCP tools",
        "Build & share the MCP server with your team",
    ], GREEN),
    ("MONTH 1", [
        "Add 2-3 more MCP servers (DB, Slack, CI/CD)",
        "Combine them: multi-tool agent workflows",
        "Document & share patterns with the org",
    ], ORANGE),
]

y_off = Inches(2.2)
for phase, items, clr in plan:
    _pill(s, Inches(0.6), y_off, Inches(1.6), Inches(0.35), phase, clr, BG, 11)
    for j, item in enumerate(items):
        _t(s, Inches(2.4), y_off + Inches(j * 0.33), Inches(4.0), Inches(0.3), item, sz=11, c=LG)
    y_off += Inches(0.33 * len(items) + 0.25)

# Quickstart
_rect(s, Inches(0.6), Inches(6.0), Inches(5.8), Inches(0.95), RGBColor(0x08, 0x0C, 0x18))
_t(s, Inches(0.8), Inches(6.05), Inches(5.4), Inches(0.25), "QUICKSTART", sz=11, c=GREEN, b=True)
_code(s, Inches(0.8), Inches(6.35), Inches(5.4), Inches(0.5),
      "pip install mcp\npython -c \"from mcp.server.fastmcp import FastMCP; print('Ready!')\"", sz=10)

# Right: Resources
_rect(s, Inches(7.0), Inches(1.6), Inches(6.0), Inches(5.6), CARD)
_bar(s, Inches(7.0), Inches(1.6), Inches(6.0), CYAN)
_t(s, Inches(7.2), Inches(1.7), Inches(5.5), Inches(0.35), "RESOURCES", sz=18, c=CYAN, b=True)

resources = [
    ("MCP Official Docs", "modelcontextprotocol.io", GREEN),
    ("MCP Python SDK", "github.com/modelcontextprotocol/python-sdk", GREEN),
    ("MCP TypeScript SDK", "github.com/modelcontextprotocol/typescript-sdk", GREEN),
    ("1000+ Pre-built Servers", "github.com/modelcontextprotocol/servers", ORANGE),
    ("VS Code MCP Setup", "code.visualstudio.com/docs/copilot/chat/mcp-servers", CYAN),
    ("Copilot Agent Mode", "github.blog/ai-and-ml/github-copilot", CYAN),
    ("OpenAI Function Calling", "platform.openai.com/docs/guides/function-calling", PURPLE),
    ("FastMCP Quickstart", "gofastmcp.com", GREEN),
]
for i, (name, url, clr) in enumerate(resources):
    y = Inches(2.2) + Inches(i * 0.55)
    _t(s, Inches(7.4), y, Inches(2.5), Inches(0.3), name, sz=12, c=clr, b=True)
    _t(s, Inches(10.0), y, Inches(2.8), Inches(0.3), url, sz=10, c=MG)

_rect(s, Inches(7.2), Inches(6.6), Inches(5.5), Inches(0.5), CARD2)
_t(s, Inches(7.2), Inches(6.62), Inches(5.5), Inches(0.45),
   "Check existing MCP servers before building from scratch!",
   sz=12, c=YELLOW, b=True, a=PP_ALIGN.CENTER)


# -- Save --
out = r"c:\MyProjects\AWS\Prompt2TestUI\AI_Agents_MCP_Agentic_Workflows.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
